# # 이미지 크기는 33에 맞춰서

# # 그냥 33에 맞춰서
# # 정방형 = 33 x 33
# # 3대4 = 33 x 44
# # 9대16 = 33 x 58

from logging import PlaceHolder
from xml.etree.ElementInclude import include
from pptx import Presentation
from pptx.util import Inches

import os

# ------------------
def printImages(fileList):
    for name in fileList:
        print(name)

# ------------------

prs = Presentation()

prs.slide_width = Inches(11.69)
prs.slide_height = Inches(8.27)

slide = prs.slides.add_slide(prs.slide_layouts[6]) # 슬라이드 추가

# 어차피 프린트를 하면 공백 있이 출력되기 때문에,
# 사진을 끝에 추가시킨다

print("사진이 있는 폴더 경로를 입력하세요 : ",end=" ")
path = input()
fileList = os.listdir(path)

printImages(fileList)

# 다음에 추가할 위치 단위는 인치
pointerLeft = 1
pointerTop = 0

imagese = []


for idx,name in enumerate(fileList):
    # left, top, 너비, 높이
    pic = slide.shapes.add_picture(path+"\\"+name,0,0,Inches(1.3)) # 3.3
    imagese.append(pic)
    width,height = pic.image.size

    print("widthPX",width)
    print("heightPX",height)

    # 96픽셀 = 1인치
    # 1픽셀 = 0.010417인치

    # 만약 w 보다 h 가 더 크다면 돌려야 한다
    if (height > width):
        pic.rotation = 90 # 돌리기
        # 가로 세로를 변경해줘야 다음 위치를 정하기 쉽다
        width,height = height,width
    else: # 사진의 가로가 더 길다면 세로를 3.3cm로 하고 가로를 다시 맞춰준다
        pic.height = Inches(1.3)
        pic.width = round((Inches(1.3) * (width / height)))


    # pic.left = Inches(0.2)
    # pic.top = Inches(-0.2)
    # if (height > width):
    #     pic.rotation = 90
    #     # 가로 세로 변경
    # else: # 만약 사진이 가로가 더 길다면 세로를 1.3인치로 해야한다
    #     pic.height = Inches(1.3)
    #     # 비율에 맞게 사진 비율을 변경시켜야 한다
    #     pic.width = Inches()


# C:\Users\DGSW\Desktop\다이어리사진

prs.save('다이어리사진.pptx')


