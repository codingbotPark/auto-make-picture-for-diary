# # 이미지 크기는 33에 맞춰서

# # 그냥 33에 맞춰서
# # 정방형 = 33 x 33
# # 3대4 = 33 x 44
# # 9대16 = 33 x 58

from logging import PlaceHolder
from xml.etree.ElementInclude import include
from pptx import Presentation
from pptx.util import Inches,Cm

from PIL import Image

import os


prs = Presentation()

prs.slide_height = Cm(21)
prs.slide_width = Cm(29.7)

slide = prs.slides.add_slide(prs.slide_layouts[6]) # 슬라이드 추가

# 어차피 프린트를 하면 공백 있이 출력되기 때문에,
# 사진을 끝에 추가시킨다

print("사진이 있는 폴더 경로를 입력하세요 : ",end=" ")
path = input()
fileList = os.listdir(path)



for imagename in fileList:
    print(imagename)
    try:
        im = Image.open(path +"\\"+ imagename)
        im.save(path +"\\" + imagename)
    except:
        print(imagename, "안됨")
        pass


# 다음에 추가할 위치 단위는 cm
pointerLeft = 0
pointerTop = 0



for idx,name in enumerate(fileList):
    # left, top, 너비, 높이
    pic = slide.shapes.add_picture(path+"\\"+name,Cm(pointerLeft),Cm(pointerTop),Cm(3.3)) # 3.3
    width,height = pic.image.size



    isOver = pointerLeft
    if (height > width): # 세로가 더 클 땐 height를 추가
        isOver += 3.3 * (height / width)
    else: # width를 추가
        isOver += 3.3 * (width / height)

    # 슬라이드 끝을 넘으면
    if (isOver > 29.7):
        pointerLeft = 0
        pointerTop += 3.3

        pic.left = Cm(pointerLeft)
        pic.top = Cm(pointerTop)

        # 슬라이드를 다 쓰면
        if pointerTop+3.3 > 21:
            pointerTop = 0
            pointerLeft = 0

            pic = pic._element
            pic.getparent().remove(pic)

            slide = prs.slides.add_slide(prs.slide_layouts[6]) # 슬라이드 추가
            pic = slide.shapes.add_picture(path+"\\"+name,Cm(pointerLeft),Cm(pointerTop),Cm(3.3)) # 3.3

    heightCM = 0
    widthCM = 0

    # 만약 w 보다 h 가 더 크다면 돌려야 한다
    if (height > width):
        pic.rotation = 90 # 돌리기
        width,height = height,width

        balance = width / height

        # 픽셀로 계산하면 안 맞기 때문에 height가 3.3cm 인 것을 활용
        heightCM = 3.3
        widthCM = round(balance * heightCM,1)

        gap = round(widthCM - heightCM,1) # 돌려서 난 차이

        # 좌우로 차이가 나기 때문에 /2를 한다
        pic.top = Cm((pic.top / 360000) - (gap/2)) 
        pic.left = Cm((pic.left / 360000) + (gap/2))


    else: # 사진의 가로가 더 길다면 세로를 3.3cm로 하고 가로를 다시 맞춰준다
        heightCM = 3.3
        widthCM = round(3.3 * (width / height),1)
        pic.height = Cm(heightCM)
        pic.width = Cm(widthCM)

    pointerLeft = widthCM + pointerLeft
    
# C:\Users\DGSW\Desktop\다이어리사진

prs.save('다이어리사진.pptx')


